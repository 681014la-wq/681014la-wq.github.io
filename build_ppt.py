"""
JooN's Sushi 메뉴 PPT v4 - AI 배경 이미지 + 프리미엄 디자인
"""
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from menu_utils import find_local_image, truncate_text

MENU_JSON = "menu.json"
IMAGE_DIR = "images"
OUTPUT_PPT = "JooNs_Sushi_Menu_v4.pptx"

# 배경 이미지
BG_COVER = "bg_cover.png"
BG_TITLE = "bg_title.png"
BG_MENU = "bg_menu.png"

# 색상
GOLD = RGBColor(0xD4, 0xAA, 0x6D)
GOLD_L = RGBColor(0xF0, 0xD4, 0xA8)
WHITE = RGBColor(0xF5, 0xF0, 0xE8)
GRAY = RGBColor(0xAA, 0xA5, 0x9E)
CARD_BG = RGBColor(0x12, 0x12, 0x1E)
CARD_BORDER = RGBColor(0x3A, 0x34, 0x28)






def load_menu():
    with open(MENU_JSON, "r", encoding="utf-8") as f:
        menu = json.load(f)
    result = []
    for item in menu:
        local = find_local_image(item["name"])
        if local:
            item["local_image"] = local
            result.append(item)
    print(f"전체: {len(menu)}개 → 이미지: {len(result)}개")
    return result


def add_bg_image(slide, prs, img_path):
    """전체 슬라이드에 배경 이미지 (없으면 단색 배경)"""
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                 prs.slide_width, prs.slide_height)
    else:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x12)


def add_text(slide, left, top, w, h, text, sz=18, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font="Segoe UI"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def add_shape(slide, x, y, w, h, fill_color, border_color=None, border_w=0, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_w)
    else:
        shape.line.fill.background()
    return shape


def make_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_image(s, prs, BG_COVER)

    add_text(s, Inches(1), Inches(2.2), Inches(11.333), Inches(1.5),
             "JooN's Sushi", sz=78, color=GOLD_L, bold=True,
             align=PP_ALIGN.CENTER, font="Georgia")

    add_text(s, Inches(1), Inches(3.9), Inches(11.333), Inches(0.6),
             "—  M  E  N  U  —", sz=30, color=WHITE,
             align=PP_ALIGN.CENTER, font="Segoe UI Light")

    add_text(s, Inches(1), Inches(5.2), Inches(11.333), Inches(0.4),
             "29910 Murrieta Hot Springs Rd L, Murrieta, CA", sz=14,
             color=GRAY, align=PP_ALIGN.CENTER, font="Segoe UI Light")

    add_text(s, Inches(1), Inches(5.6), Inches(11.333), Inches(0.3),
             "joonssushimurrieta.com", sz=12,
             color=GOLD, align=PP_ALIGN.CENTER)


def make_category_title(prs, cat_name, count):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_image(s, prs, BG_TITLE)

    add_text(s, Inches(1), Inches(2.8), Inches(11.333), Inches(1.2),
             cat_name, sz=56, color=GOLD_L, bold=True,
             align=PP_ALIGN.CENTER, font="Georgia")

    add_text(s, Inches(1), Inches(4.2), Inches(11.333), Inches(0.4),
             f"✦  {count} Selections  ✦", sz=16, color=GRAY,
             align=PP_ALIGN.CENTER, font="Segoe UI Light")


def make_menu_slide(prs, items, cat_name):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_image(s, prs, BG_MENU)

    n = len(items)

    # 카테고리명 (상단)
    add_text(s, Inches(0.5), Inches(0.15), Inches(8), Inches(0.45),
             cat_name, sz=18, color=GOLD_L, bold=True, font="Georgia")

    add_text(s, Inches(9.5), Inches(0.18), Inches(3.5), Inches(0.4),
             "JooN's Sushi", sz=11, color=GOLD,
             align=PP_ALIGN.RIGHT, font="Georgia")

    # 개수별 레이아웃
    if n == 1:
        card_w = 6.5
        img_sz = 4.8
        name_sz = 28
        price_sz = 36
        desc_sz = 14
        gap = 0
    elif n == 2:
        card_w = 5.4
        img_sz = 4.0
        name_sz = 22
        price_sz = 30
        desc_sz = 12
        gap = 0.6
    else:
        card_w = 3.8
        img_sz = 3.1
        name_sz = 16
        price_sz = 24
        desc_sz = 10
        gap = 0.35

    total_w = n * card_w + (n - 1) * gap
    sx = (13.333 - total_w) / 2
    card_h = 6.3

    for i, item in enumerate(items):
        x = Inches(sx + i * (card_w + gap))
        y = Inches(0.8)

        # 카드 배경 (반투명 다크)
        card = add_shape(s, x, y, Inches(card_w), Inches(card_h),
                         CARD_BG, CARD_BORDER, 1)

        # 카드 상단 골드 라인
        add_shape(s, x, y, Inches(card_w), Pt(3), GOLD)

        # 이미지
        img_m = (card_w - img_sz) / 2
        img_x = x + Inches(img_m)
        img_y = y + Inches(0.3)

        if item["local_image"] and os.path.exists(item["local_image"]):
            try:
                s.shapes.add_picture(item["local_image"],
                                     img_x, img_y,
                                     Inches(img_sz), Inches(img_sz))
            except Exception:
                pass

        # 이미지 아래 골드 장식 라인
        deco_y = img_y + Inches(img_sz) + Inches(0.08)
        line_w = img_sz * 0.5
        line_x = img_x + Inches((img_sz - line_w) / 2)
        add_shape(s, line_x, deco_y, Inches(line_w), Pt(1.5), GOLD)

        # 이름
        name_y = deco_y + Inches(0.15)
        add_text(s, x + Inches(0.2), name_y,
                 Inches(card_w - 0.4), Inches(0.8),
                 item["name"], sz=name_sz, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, font="Georgia")

        # 가격
        price_y = name_y + Inches(0.6) if n == 3 else name_y + Inches(0.7)
        add_text(s, x + Inches(0.2), price_y,
                 Inches(card_w - 0.4), Inches(0.5),
                 item["price"], sz=price_sz, color=GOLD_L, bold=True,
                 align=PP_ALIGN.CENTER, font="Georgia")

        # 설명
        desc = item.get("description", "")
        if desc:
            desc_y = price_y + Inches(0.5) if n == 3 else price_y + Inches(0.6)
            desc = truncate_text(desc, 80)
            add_text(s, x + Inches(0.2), desc_y,
                     Inches(card_w - 0.4), Inches(1.0),
                     desc, sz=desc_sz, color=GRAY,
                     align=PP_ALIGN.CENTER, font="Segoe UI Light")


def create_ppt(menu):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_cover(prs)

    cats = {}
    for item in menu:
        c = item["category"]
        if c not in cats:
            cats[c] = []
        cats[c].append(item)

    for cat, items in cats.items():
        make_category_title(prs, cat, len(items))
        for si in range(0, len(items), 3):
            chunk = items[si:si + 3]
            make_menu_slide(prs, chunk, cat)

    prs.save(OUTPUT_PPT)
    print(f"\nPPT 저장 완료: {OUTPUT_PPT}")
    print(f"   슬라이드: {len(prs.slides)}장")


if __name__ == "__main__":
    menu = load_menu()
    create_ppt(menu)
